import os
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from datetime import datetime
from flask import Blueprint, request, jsonify
from werkzeug.utils import secure_filename
from PIL import Image
import pytesseract
from db.models import db, AcademicResource
from search.embedding_service import generate_embedding
from search.faiss_store import add_embedding, find_duplicates
from classification.category_service import classify_text
from db.activity_logger import log_activity

ingestion_bp = Blueprint("ingestion", __name__)

UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

IMAGE_EXTENSIONS = {"png", "jpg", "jpeg"}
PDF_EXTENSIONS = {"pdf"}
DOCX_EXTENSIONS = {"docx"}
PPTX_EXTENSIONS = {"pptx"}
OCR_CONFIDENCE_THRESHOLD = 60  # out of 100, per FR3

def get_extension(filename):
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


@ingestion_bp.route("/resources/upload", methods=["POST"])
def upload_resource():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files["file"]
    user_id = request.form.get("userID")

    if not user_id:
        return jsonify({"error": "userID is required"}), 400
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    filename = secure_filename(file.filename)
    extension = get_extension(filename)
    save_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(save_path)

    extracted_text = None
    confidence = None
    needs_review = False

    if extension in IMAGE_EXTENSIONS:
        try:
            image = Image.open(save_path)
            ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            words = [w for w in ocr_data["text"] if w.strip() != ""]
            confidences = [int(c) for c, w in zip(ocr_data["conf"], ocr_data["text"]) if w.strip() != "" and c != "-1"]

            extracted_text = " ".join(words)
            confidence = sum(confidences) / len(confidences) if confidences else 0
            needs_review = confidence < OCR_CONFIDENCE_THRESHOLD

        except Exception as e:
            return jsonify({"error": f"OCR failed: {str(e)}"}), 500

    elif extension in PDF_EXTENSIONS:
     try:
            text_parts = []
            with pdfplumber.open(save_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            extracted_text = "\n".join(text_parts).strip()

            if extracted_text:
                confidence = 100.0  # digitally-extracted text, not OCR'd — treat as fully trustworthy
                needs_review = False
            else:
                # scanned/image-only PDF with no extractable text layer
                extracted_text = None
                confidence = 0.0
                needs_review = True

     except Exception as e:
            return jsonify({"error": f"PDF text extraction failed: {str(e)}"}), 500

    elif extension in DOCX_EXTENSIONS:
        try:
            doc = DocxDocument(save_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            extracted_text = "\n".join(text_parts).strip()
            confidence = 100.0 if extracted_text else 0.0
            needs_review = not bool(extracted_text)
        except Exception as e:
            return jsonify({"error": f"DOCX text extraction failed: {str(e)}"}), 500

    elif extension in PPTX_EXTENSIONS:
        try:
            prs = Presentation(save_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            extracted_text = "\n".join(text_parts).strip()
            confidence = 100.0 if extracted_text else 0.0
            needs_review = not bool(extracted_text)
        except Exception as e:
            return jsonify({"error": f"PPTX text extraction failed: {str(e)}"}), 500

    resource = AcademicResource(
        userID=user_id,
        title=filename,
        fileName=filename,
        filePath=save_path,
        fileType=extension,
        extractedText=extracted_text,
        ocrConfidence=confidence,
        needsReview=needs_review,
    )
    db.session.add(resource)
    db.session.commit()

    embedding_vector_id = None
    duplicate_warning = None
    category_name = None
    category_confidence = None

    if extracted_text:
        embedding = generate_embedding(extracted_text)
        if embedding is not None:
            duplicates = find_duplicates(embedding, exclude_resource_id=resource.resourceID)
            if duplicates:
                duplicate_warning = duplicates

            embedding_vector_id = add_embedding(resource.resourceID, embedding)
            resource.embeddingVectorID = embedding_vector_id

            category_id, category_name, category_confidence = classify_text(embedding)
            resource.categoryID = category_id

            db.session.commit()

        log_activity(user_id, "upload", resource.resourceID)

    return jsonify({
        "message": "Resource uploaded successfully",
        "resourceID": resource.resourceID,
        "extractedText": extracted_text,
        "ocrConfidence": confidence,
        "needsReview": needs_review,
        "duplicateWarning": duplicate_warning,
        "category": category_name,
        "categoryConfidence": round(category_confidence * 100, 1) if category_confidence is not None else None,
    }), 201